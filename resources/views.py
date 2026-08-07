# resources/views.py
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
from django.contrib import messages
from django.db.models import Q, F, Count, Sum, Avg
from django.contrib.auth.forms import PasswordResetForm
from django.core.paginator import Paginator
from django.utils import timezone
from django.contrib.auth.models import User
from django.conf import settings
from django.db.models.functions import TruncDate, Coalesce
import json
import zipfile
import requests
from io import BytesIO

from accounts.models import Profile
from .models import (
    Resource,
    SEMESTER_CHOICES,
    RESOURCE_TYPE_CHOICES,
    Subject,
    Favorite,
    Comment,
    Rating,
    Notification,
    Report,
    Visit,
)

from .forms import ResourceForm, CommentForm, RatingForm, ReportForm, SubjectForm, AdminResourceForm
from core import ai_service

# extra optional preview libs
try:
    import docx  # for DOCX preview
except ImportError:
    docx = None

try:
    import pptx  # for PPT/PPTX preview
except ImportError:
    pptx = None


# -------------------------------------------------------------------
# Resource list + filters
# -------------------------------------------------------------------
@login_required
def resource_list(request):
    # Optimize N+1 queries by selecting related owner and subject
    resources = Resource.objects.select_related("owner", "subject").all().order_by("-created_at")

    q = request.GET.get("q", "").strip()
    # Use getlist for multiple selections
    subjects = request.GET.getlist("subject")
    semesters = request.GET.getlist("semester")
    resource_types = request.GET.getlist("resource_type")
    sort = request.GET.get("sort", "newest")

    # Clean up empty strings from lists if any
    subjects = [s for s in subjects if s]
    semesters = [s for s in semesters if s]
    resource_types = [rt for rt in resource_types if rt]

    # Search (combine tokens)
    if q:
        tokens = q.split()
        for token in tokens:
            resources = resources.filter(
                Q(title__icontains=token)
                | Q(description__icontains=token)
                | Q(subject__name__icontains=token)
                | Q(subject__abbreviation__icontains=token)   # match "OS", "DBMS", etc.
                | Q(owner__username__icontains=token)
            )

    # Subject filter
    if subjects:
        try:
            subject_ids = [int(s) for s in subjects]
            resources = resources.filter(subject_id__in=subject_ids)
        except ValueError:
            pass

    # Semester filter
    if semesters:
        try:
            sem_values = [int(s) for s in semesters]
            resources = resources.filter(semester__in=sem_values)
        except ValueError:
            pass

    # Resource type filter
    if resource_types:
        resources = resources.filter(resource_type__in=resource_types)

    # Sorting
    if sort == "oldest":
        resources = resources.order_by("created_at")
    elif sort == "downloads":
        resources = resources.order_by("-download_count", "-created_at")
    elif sort == "rating":
        resources = resources.annotate(
            avg_rating=Avg("ratings__stars")
        ).order_by("-avg_rating", "-created_at")
    elif sort == "az":
        resources = resources.order_by("title")
    elif sort == "subject":
        resources = resources.order_by("subject__name", "title")
    elif sort == "ai_rank":
        resources = resources.annotate(
            ai_score=(F("download_count") * 0.5) + (F("view_count") * 0.1)
        ).order_by("-ai_score")
    else:  # newest
        resources = resources.order_by("-created_at")

    paginator = Paginator(resources, 10)
    page_number = request.GET.get("page")
    page_obj = paginator.get_page(page_number)

    top_downloads = Resource.objects.select_related("owner", "subject").order_by("-download_count").first()

    query_params = request.GET.copy()
    if "page" in query_params:
        del query_params["page"]
    query_string = query_params.urlencode()

    context = {
        "page_obj": page_obj,
        "q": q,
        "subject_filter": subjects,
        "semester_filter": semesters,
        "resource_type_filter": resource_types,
        "sort": sort,
        "semester_choices": SEMESTER_CHOICES,
        "resource_type_choices": RESOURCE_TYPE_CHOICES,
        "subjects": Subject.objects.order_by("name"),
        "top_downloads": top_downloads,
        "query_string": query_string,
    }
    return render(request, "resources/resource_list.html", context)


# -------------------------------------------------------------------
# Upload resource
# -------------------------------------------------------------------
from core import ai_service

@login_required
def upload_resource(request):
    if request.method == "POST":
        form = ResourceForm(request.POST, request.FILES)
        if form.is_valid():
            resource = form.save(commit=False)
            resource.owner = request.user

            # Handle new subject
            if not resource.subject and form.cleaned_data.get("new_subject_name"):
                new_sub_name = form.cleaned_data["new_subject_name"].strip()
                if new_sub_name:
                    new_sub, created = Subject.objects.get_or_create(name=new_sub_name)
                    resource.subject = new_sub

            resource.save()
            
            # Extract text for AI if possible
            extracted_text = ""
            ext = (resource.file_ext or "").lower()
            text_pages_dict = {}
            
            if ext == "pdf":
                try:
                    import pypdf
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    reader = pypdf.PdfReader(BytesIO(response.content))
                    
                    text_pages = []
                    # Extract from first 30 pages max to save time/memory but provide decent context
                    for i, page in enumerate(reader.pages[:30]):
                        page_text = page.extract_text() or ""
                        text_pages.append(page_text)
                        text_pages_dict[str(i + 1)] = page_text
                    extracted_text = "\n".join(text_pages)
                except Exception:
                    pass
            elif ext == "docx" and docx:
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    document = docx.Document(BytesIO(response.content))
                    extracted_text = "\n".join([p.text for p in document.paragraphs if p.text.strip()][:50])
                    if extracted_text:
                        text_pages_dict = {"1": extracted_text}
                except Exception:
                    pass
            elif ext in ["ppt", "pptx"] and pptx:
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    prs = pptx.Presentation(BytesIO(response.content))
                    all_text_parts = []
                    for i, slide in enumerate(prs.slides, start=1):
                        slide_lines = []
                        # Title first
                        if slide.shapes.title and slide.shapes.title.text.strip():
                            slide_lines.append(slide.shapes.title.text.strip())
                        # All other text shapes (body, bullet points, text boxes)
                        for shape in slide.shapes:
                            if shape == slide.shapes.title:
                                continue
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    line = para.text.strip()
                                    if line:
                                        slide_lines.append(line)
                        slide_text = "\n".join(slide_lines)
                        if slide_text.strip():
                            all_text_parts.append(slide_text)
                            # Each slide = one "page" for AI chat
                            text_pages_dict[str(i)] = slide_text
                    extracted_text = "\n\n".join(all_text_parts)
                except Exception:
                    pass
            elif ext == "txt":
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    extracted_text = response.text[:2000]
                    if extracted_text:
                        text_pages_dict = {"1": extracted_text}
                except Exception:
                    pass
                    
            if not extracted_text.strip():
                extracted_text = resource.title + "\n" + resource.description
                text_pages_dict = {"1": extracted_text}
            
            # Save the page mapping for AI Chat
            resource.extracted_text_pages = text_pages_dict

            # Generate AI fields
            resource.auto_summary = ai_service.generate_summary(extracted_text)
            resource.auto_questions = ai_service.generate_important_questions(extracted_text)
            resource.save()

            messages.success(request, "Resource uploaded successfully! AI has processed it.")
            return redirect("resources:resource_list")
        else:
            messages.error(request, "Please correct the errors in the form.")
    else:
        form = ResourceForm()

    return render(request, "resources/upload_resource.html", {"form": form})


# -------------------------------------------------------------------
# Resource detail (comments + ratings + favorite)
# -------------------------------------------------------------------
@login_required
def resource_detail(request, pk):
    resource = get_object_or_404(Resource.objects.select_related("owner", "subject"), pk=pk)

    # Count views only on GET
    if request.method == "GET":
        Resource.objects.filter(pk=pk).update(view_count=F("view_count") + 1)
        resource.refresh_from_db()

    comments = (
        resource.comments.select_related("user")
        .prefetch_related("replies__user")
        .filter(parent__isnull=True)
    )

    if request.method == "POST":
        # ----- New comment / reply -----
        if "comment_submit" in request.POST:
            comment_form = CommentForm(request.POST)
            rating_form = RatingForm()  # untouched
            if comment_form.is_valid():
                parent = None
                parent_id = request.POST.get("parent_id")
                if parent_id:
                    parent = Comment.objects.filter(
                        pk=parent_id, resource=resource
                    ).first()

                comment = comment_form.save(commit=False)
                comment.resource = resource
                comment.user = request.user
                comment.parent = parent
                comment.save()

                # Notify owners
                if resource.owner != request.user:
                    Notification.objects.create(
                        user=resource.owner,
                        notif_type="COMMENT" if parent is None else "REPLY",
                        message=(
                            f"{request.user.username} "
                            f"{'commented on' if parent is None else 'replied on'} "
                            f"your resource '{resource.title}'."
                        ),
                        resource=resource,
                        comment=comment,
                    )

                # Notify parent commenter if different
                if parent and parent.user not in (request.user, resource.owner):
                    Notification.objects.create(
                        user=parent.user,
                        notif_type="REPLY",
                        message=(
                            f"{request.user.username} replied to your comment "
                            f"on '{resource.title}'."
                        ),
                        resource=resource,
                        comment=comment,
                    )

                messages.success(request, "Comment added!")
                return redirect("resources:resource_detail", pk=resource.pk)

        # ----- Rating submit -----
        elif "rating_submit" in request.POST:
            rating_form = RatingForm(request.POST)
            comment_form = CommentForm()
            if rating_form.is_valid():
                stars = rating_form.cleaned_data["stars"]
                Rating.objects.update_or_create(
                    resource=resource,
                    user=request.user,
                    defaults={"stars": stars},
                )

                if resource.owner != request.user:
                    Notification.objects.create(
                        user=resource.owner,
                        notif_type="RATING",
                        message=(
                            f"{request.user.username} rated your resource "
                            f"'{resource.title}' with {stars}★."
                        ),
                        resource=resource,
                    )

                messages.success(request, "Your rating has been saved!")
                return redirect("resources:resource_detail", pk=resource.pk)
    else:
        comment_form = CommentForm()
        try:
            existing = Rating.objects.get(resource=resource, user=request.user)
            rating_form = RatingForm(initial={"stars": existing.stars})
        except Rating.DoesNotExist:
            rating_form = RatingForm()

    is_favorite = resource.is_favorited_by(request.user)

    # AI recommendations
    recommendations = ai_service.get_recommendations(resource)

    context = {
        "resource": resource,
        "resource_type_choices": RESOURCE_TYPE_CHOICES,
        "comment_form": comment_form,
        "rating_form": rating_form,
        "comments": comments,
        "is_favorite": is_favorite,
        "recommendations": recommendations,
    }
    return render(request, "resources/resource_detail.html", context)


# -------------------------------------------------------------------
# AI Chat API
# -------------------------------------------------------------------
@login_required
def resource_chat_api(request, pk):
    """
    API endpoint to handle Chat with PDF messages.
    """
    import json
    from django.http import JsonResponse
    from core import ai_service
    from django.views.decorators.csrf import csrf_exempt
    
    if request.method != "POST":
        return JsonResponse({"error": "Only POST requests allowed."}, status=405)
        
    resource = get_object_or_404(Resource, pk=pk)
    
    try:
        data = json.loads(request.body)
        chat_history = data.get("history", [])
        new_question = data.get("message", "").strip()
        
        if not new_question:
            return JsonResponse({"error": "Empty message"}, status=400)

        # ---- Use stored extraction OR live-extract if missing/stale ----
        pages_dict = resource.extracted_text_pages or {}

        # Check if stored data is stale (only 1 page of < 200 chars = old title-only extraction)
        is_stale = (
            not pages_dict
            or (len(pages_dict) == 1 and len(list(pages_dict.values())[0]) < 200)
        )

        if is_stale:
            ext = (resource.file_ext or "").lower()

            if ext in ["ppt", "pptx"] and pptx:
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    prs = pptx.Presentation(BytesIO(response.content))
                    fresh_pages = {}
                    for i, slide in enumerate(prs.slides, start=1):
                        slide_lines = []
                        if slide.shapes.title and slide.shapes.title.text.strip():
                            slide_lines.append(slide.shapes.title.text.strip())
                        for shape in slide.shapes:
                            if shape == slide.shapes.title:
                                continue
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    line = para.text.strip()
                                    if line:
                                        slide_lines.append(line)
                        slide_text = "\n".join(slide_lines)
                        if slide_text.strip():
                            fresh_pages[str(i)] = slide_text
                    if fresh_pages:
                        pages_dict = fresh_pages
                        # Save back so next chat is instant
                        Resource.objects.filter(pk=pk).update(extracted_text_pages=fresh_pages)
                except Exception:
                    pass

            elif ext == "pdf":
                try:
                    import pypdf
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    reader = pypdf.PdfReader(BytesIO(response.content))
                    fresh_pages = {}
                    for i, page in enumerate(reader.pages[:30], start=1):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            fresh_pages[str(i)] = page_text
                    if fresh_pages:
                        pages_dict = fresh_pages
                        Resource.objects.filter(pk=pk).update(extracted_text_pages=fresh_pages)
                except Exception:
                    pass

            elif ext == "docx" and docx:
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    document = docx.Document(BytesIO(response.content))
                    full_text = "\n".join(
                        [p.text.strip() for p in document.paragraphs if p.text.strip()]
                    )
                    if full_text:
                        fresh_pages = {"1": full_text}
                        pages_dict = fresh_pages
                        Resource.objects.filter(pk=pk).update(extracted_text_pages=fresh_pages)
                except Exception:
                    pass

        # Call AI Service
        answer = ai_service.chat_with_document(
            extracted_pages_dict=pages_dict,
            chat_history=chat_history,
            new_question=new_question
        )

        return JsonResponse({"answer": answer})
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


# -------------------------------------------------------------------
# Download
# -------------------------------------------------------------------
@login_required
def resource_download(request, pk):
    resource = get_object_or_404(Resource, pk=pk)
    Resource.objects.filter(pk=pk).update(download_count=F("download_count") + 1)
    return redirect(resource.file.url)


# -------------------------------------------------------------------
# ONLINE VIEWER (image / pdf / docx / ppt / zip)
# -------------------------------------------------------------------
@login_required
def resource_viewer(request, pk):
    """
    Online viewer for various file types.
    """
    resource = get_object_or_404(Resource, pk=pk)

    # Count as a view
    Resource.objects.filter(pk=pk).update(view_count=F("view_count") + 1)
    resource.refresh_from_db()

    ext = (resource.file_ext or "").lower()
    preview_type = "fallback"
    preview_data = None

    try:
        # ---------- 1) IMAGE ----------
        if resource.is_image:
            preview_type = "image"
            preview_data = resource.file.url

        # ---------- 2) PDF ----------
        elif ext == "pdf":
            preview_type = "pdf"
            preview_data = resource.file.url

        # ---------- 3) DOCX ----------
        elif ext == "docx":
            preview_type = "docx"
            if docx is None:
                preview_data = None
            else:
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    document = docx.Document(BytesIO(response.content))
                    lines = [p.text.strip() for p in document.paragraphs if p.text.strip()]
                    preview_data = lines[:40]  # first 40 lines
                except Exception:
                    preview_data = None

        # ---------- 4) PPT / PPTX ----------
        elif ext in ["ppt", "pptx"]:
            preview_type = "ppt"
            if pptx is None:
                preview_data = None
            else:
                try:
                    response = requests.get(resource.file.url)
                    response.raise_for_status()

                    prs = pptx.Presentation(BytesIO(response.content))
                    slides_data = []
                    for i, slide in enumerate(prs.slides, start=1):
                        # Get slide title
                        title_text = ""
                        if slide.shapes.title and slide.shapes.title.text.strip():
                            title_text = slide.shapes.title.text.strip()
                        else:
                            title_text = f"Slide {i}"

                        # Get all body text (bullet points, text boxes, etc.)
                        body_lines = []
                        for shape in slide.shapes:
                            # Skip the title shape to avoid duplication
                            if shape == slide.shapes.title:
                                continue
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    line = para.text.strip()
                                    if line:
                                        body_lines.append(line)

                        slides_data.append({
                            "number": i,
                            "title": title_text,
                            "body": body_lines,
                        })
                    preview_data = slides_data
                except Exception as e:
                    preview_data = None

        # ---------- 5) ZIP ----------
        elif ext == "zip":
            preview_type = "zip"
            try:
                names = []
                response = requests.get(resource.file.url)
                response.raise_for_status()

                zip_data = BytesIO(response.content)

                with zipfile.ZipFile(zip_data, "r") as zf:
                    for info in zf.infolist():
                        if not info.is_dir():
                            names.append(info.filename)
                preview_data = names
            except Exception:
                preview_data = None

        # ---------- 6) OTHER FILES ----------
        else:
            preview_type = "fallback"
            preview_data = None

    except Exception:
        preview_type = "fallback"
        preview_data = None

    return render(
        request,
        "resources/resource_viewer.html",
        {
            "resource": resource,
            "preview_type": preview_type,
            "preview_data": preview_data,
        },
    )


# -------------------------------------------------------------------
# Favorites
# -------------------------------------------------------------------
@login_required
def toggle_favorite(request, pk):
    resource = get_object_or_404(Resource, pk=pk)
    favorite, created = Favorite.objects.get_or_create(
        user=request.user, resource=resource
    )
    if created:
        messages.success(request, "Added to your favorites!")
    else:
        favorite.delete()
        messages.info(request, "Removed from your favorites.")
    return redirect("resources:resource_detail", pk=pk)


@login_required
def my_favorites(request):
    favorites = Favorite.objects.filter(user=request.user).select_related(
        "resource", "resource__subject", "resource__owner"
    )
    resources = [f.resource for f in favorites]
    return render(
        request,
        "resources/my_favorites.html",
        {
            "resources": resources,
            "resource_type_choices": RESOURCE_TYPE_CHOICES,
        },
    )

# -------------------------------------------------------------------
# Subject dashboard (for charts)
# -------------------------------------------------------------------
@login_required
def subject_dashboard(request):
    subject_stats = (
        Resource.objects.values("subject__name")
        .annotate(
            total_uploads=Count("id"),
            total_downloads=Sum("download_count"),
            total_views=Sum("view_count"),
        )
        .order_by("subject__name")
    )

    labels = []
    uploads = []
    downloads = []
    views = []

    for row in subject_stats:
        labels.append(row["subject__name"] or "N/A")
        uploads.append(row["total_uploads"] or 0)
        downloads.append(row["total_downloads"] or 0)
        views.append(row["total_views"] or 0)

    context = {
        "labels": labels,
        "uploads": uploads,
        "downloads": downloads,
        "views": views,
    }
    return render(request, "resources/subject_dashboard.html", context)


# -------------------------------------------------------------------
# Report resource
# -------------------------------------------------------------------
@login_required
def report_resource(request, pk):
    resource = get_object_or_404(Resource, pk=pk)

    if request.method == "POST":
        form = ReportForm(request.POST)
        if form.is_valid():
            report = form.save(commit=False)
            report.resource = resource
            report.reporter = request.user
            report.save()

            if resource.owner != request.user:
                Notification.objects.create(
                    user=resource.owner,
                    notif_type="REPORT",
                    message=(
                        f"Your resource '{resource.title}' was reported by "
                        f"{request.user.username}."
                    ),
                    resource=resource,
                    report=report,
                )

            messages.success(request, "Thank you. Your report has been submitted.")
        else:
            messages.error(request, "Could not submit report, please check the form.")
    return redirect("resources:resource_detail", pk=pk)


# -------------------------------------------------------------------
# My activity dashboard
# -------------------------------------------------------------------
@login_required
def my_activity(request):
    """
    Shows stats for the current user:
    - total uploads
    - favorites saved
    - ratings received
    - total views & downloads on their uploads
    - chart for views/downloads per upload
    - detailed table with each uploaded resource
    """
    user = request.user

    uploads_qs = Resource.objects.filter(owner=user).select_related("owner", "subject")

    uploads_count = uploads_qs.count()
    favorites_count = Favorite.objects.filter(user=user).count()
    ratings_received = Rating.objects.filter(resource__owner=user).count()

    total_views = uploads_qs.aggregate(total=Sum("view_count"))["total"] or 0
    total_downloads = uploads_qs.aggregate(total=Sum("download_count"))["total"] or 0

    # latest 10 uploads in chronological order (oldest first among those 10)
    uploads_for_chart = uploads_qs.order_by("-created_at")[:10][::-1]

    labels = [r.title[:20] for r in uploads_for_chart]
    views_data = [r.view_count for r in uploads_for_chart]
    downloads_data = [r.download_count for r in uploads_for_chart]

    my_uploads = uploads_qs.order_by("-created_at")

    context = {
        "uploads_count": uploads_count,
        "favorites_count": favorites_count,
        "ratings_received": ratings_received,
        "total_views": total_views,
        "total_downloads": total_downloads,
        "labels_json": json.dumps(labels),
        "views_json": json.dumps(views_data),
        "downloads_json": json.dumps(downloads_data),
        "my_uploads": my_uploads,
        "resource_type_choices": RESOURCE_TYPE_CHOICES,
    }
    return render(request, "resources/my_activity.html", context)


# -------------------------------------------------------------------
# Admin analytics dashboard
# -------------------------------------------------------------------
@staff_member_required
def admin_analytics_dashboard(request):
    """
    Analytics for staff: uploads, downloads, reports, visits, active users.
    """
    # --- Totals ---
    total_resources = Resource.objects.count()
    total_downloads = Resource.objects.aggregate(total=Sum("download_count"))["total"] or 0
    total_views = Resource.objects.aggregate(total=Sum("view_count"))["total"] or 0
    total_reports = Report.objects.count()
    total_visits = Visit.objects.count()

    # --- Top resources by downloads ---
    top_resources = Resource.objects.order_by("-download_count")[:5]

    # --- Subject-wise uploads ---
    subject_stats = (
        Resource.objects.values("subject__name")
        .annotate(count=Count("id"))
        .order_by("-count")
    )

    # --- Uploads per day ---
    uploads_by_day_qs = (
        Resource.objects
        .annotate(day=TruncDate("created_at"))
        .values("day")
        .annotate(count=Count("id"))
        .order_by("day")
    )
    uploads_labels = [row["day"].strftime("%Y-%m-%d") for row in uploads_by_day_qs]
    uploads_counts = [row["count"] for row in uploads_by_day_qs]

    # --- Visits per day ---
    visits_by_day_qs = (
        Visit.objects
        .annotate(day=TruncDate("created_at"))
        .values("day")
        .annotate(count=Count("id"))
        .order_by("day")
    )
    visits_labels = [row["day"].strftime("%Y-%m-%d") for row in visits_by_day_qs]
    visits_counts = [row["count"] for row in visits_by_day_qs]

    # --- Downloads per subject ---
    subject_download_qs = (
        Resource.objects.values("subject__name")
        .annotate(downloads=Sum("download_count"))
        .order_by("-downloads")
    )
    subject_download_labels = [
        row["subject__name"] or "N/A" for row in subject_download_qs
    ]
    subject_download_counts = [row["downloads"] or 0 for row in subject_download_qs]

    # --- Most active users (by uploads) ---
    active_users_qs = (
        User.objects.annotate(
            uploads_count=Count("resources", distinct=True),
            ratings_given=Count("ratings", distinct=True),
            comments_made=Count("comments", distinct=True),
        )
        .order_by("-uploads_count")[:5]
    )

    context = {
        "total_resources": total_resources,
        "total_downloads": total_downloads,
        "total_views": total_views,
        "total_reports": total_reports,
        "total_visits": total_visits,
        "top_resources": top_resources,
        "subject_stats": subject_stats,
        "active_users": active_users_qs,
        "uploads_labels_json": json.dumps(uploads_labels),
        "uploads_counts_json": json.dumps(uploads_counts),
        "visits_labels_json": json.dumps(visits_labels),
        "visits_counts_json": json.dumps(visits_counts),
        "subject_download_labels_json": json.dumps(subject_download_labels),
        "subject_download_counts_json": json.dumps(subject_download_counts),
    }
    return render(request, "resources/admin_analytics.html", context)


# -------------------------------------------------------------------
# Admin center helpers
# -------------------------------------------------------------------

def _get_role_label(user):
    if user.is_superuser:
        return "Admin"
    if user.is_staff:
        return "Moderator"
    return "Student"


def _require_superuser(request):
    if not request.user.is_superuser:
        messages.error(request, "Superuser access required.")
        return False
    return True


@staff_member_required
def admin_user_list(request):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    q = request.GET.get("q", "").strip()
    role = request.GET.get("role", "")
    status = request.GET.get("status", "")

    users = (
        User.objects.select_related("profile")
        .annotate(
            uploads_count=Count("resources", distinct=True),
            reports_count=Count("reports", distinct=True),
            comments_count=Count("comments", distinct=True),
        )
        .order_by("-date_joined")
    )

    if q:
        users = users.filter(
            Q(username__icontains=q)
            | Q(email__icontains=q)
            | Q(profile__college__icontains=q)
            | Q(profile__branch__icontains=q)
        )

    if role == "student":
        users = users.filter(is_staff=False, is_superuser=False)
    elif role == "moderator":
        users = users.filter(is_staff=True, is_superuser=False)
    elif role == "admin":
        users = users.filter(is_superuser=True)

    if status == "active":
        users = users.filter(is_active=True)
    elif status == "inactive":
        users = users.filter(is_active=False)

    paginator = Paginator(users, 20)
    page_number = request.GET.get("page")
    page_obj = paginator.get_page(page_number)

    for user in page_obj:
        user.role_label = _get_role_label(user)

    context = {
        "page_obj": page_obj,
        "q": q,
        "role": role,
        "status": status,
        "Profile": Profile,
    }
    return render(request, "resources/admin_user_list.html", context)


@staff_member_required
def admin_user_detail(request, user_id):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    target_user = get_object_or_404(User, pk=user_id)
    profile, _ = Profile.objects.get_or_create(user=target_user)

    if request.method == "POST":
        action = request.POST.get("action")

        if action == "toggle_active":
            target_user.is_active = not target_user.is_active
            target_user.save()
            messages.success(request, "User status updated.")
        elif action == "change_role":
            role = request.POST.get("role")
            if role == "admin":
                target_user.is_staff = True
                target_user.is_superuser = True
            elif role == "moderator":
                target_user.is_staff = True
                target_user.is_superuser = False
            else:
                target_user.is_staff = False
                target_user.is_superuser = False
            target_user.save()
            messages.success(request, "User role updated.")
        elif action == "verify_email":
            profile.email_verified = True
            profile.email_verified_at = timezone.now()
            profile.save()
            messages.success(request, "User email verified manually.")
        elif action == "send_reset":
            if target_user.email:
                form = PasswordResetForm({"email": target_user.email})
                if form.is_valid():
                    form.save(
                        request=request,
                        use_https=request.is_secure(),
                        subject_template_name="accounts/password_reset_subject.txt",
                        email_template_name="accounts/password_reset_email.html",
                        from_email=getattr(settings, "DEFAULT_FROM_EMAIL", None),
                        html_email_template_name=None,
                    )
                    messages.success(request, "Password reset email sent to user.")
                else:
                    messages.error(request, "Unable to send reset email for this account.")
            else:
                messages.error(request, "User has no email address configured.")
        return redirect("resources:admin_user_detail", user_id=target_user.id)

    user_stats = {
        "uploads_count": target_user.resources.count(),
        "favorites_count": target_user.favorites.count(),
        "reports_count": target_user.reports.count(),
        "comments_count": target_user.comments.count(),
        "ratings_given": target_user.ratings.count(),
    }

    recent_uploads = target_user.resources.order_by("-created_at")[:5]

    context = {
        "target_user": target_user,
        "profile": profile,
        "user_stats": user_stats,
        "recent_uploads": recent_uploads,
        "role_label": _get_role_label(target_user),
    }
    return render(request, "resources/admin_user_detail.html", context)


@staff_member_required
def admin_resource_list(request):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    resources_qs = Resource.objects.select_related("owner", "subject").all().order_by("-created_at")
    q = request.GET.get("q", "").strip()
    subject_ids = request.GET.getlist("subject")
    semesters = request.GET.getlist("semester")
    status = request.GET.get("status", "")
    sort = request.GET.get("sort", "newest")

    if q:
        tokens = q.split()
        for token in tokens:
            resources_qs = resources_qs.filter(
                Q(title__icontains=token)
                | Q(description__icontains=token)
                | Q(subject__name__icontains=token)
                | Q(subject__abbreviation__icontains=token)
                | Q(owner__username__icontains=token)
            )

    if subject_ids:
        try:
            resources_qs = resources_qs.filter(subject_id__in=[int(s) for s in subject_ids if s])
        except ValueError:
            pass

    if semesters:
        try:
            resources_qs = resources_qs.filter(semester__in=[int(s) for s in semesters if s])
        except ValueError:
            pass

    if status == "pending":
        resources_qs = resources_qs.filter(verification_status="PENDING")
    elif status == "approved":
        resources_qs = resources_qs.filter(verification_status="APPROVED")
    elif status == "rejected":
        resources_qs = resources_qs.filter(verification_status="REJECTED")

    if sort == "downloads":
        resources_qs = resources_qs.order_by("-download_count", "-created_at")
    elif sort == "views":
        resources_qs = resources_qs.order_by("-view_count", "-created_at")
    elif sort == "title":
        resources_qs = resources_qs.order_by("title")
    else:
        resources_qs = resources_qs.order_by("-created_at")

    if request.method == "POST" and request.POST.get("bulk_action"):
        action = request.POST.get("bulk_action")
        selected_ids = request.POST.getlist("selected_resources")
        selected_resources = Resource.objects.filter(pk__in=selected_ids)

        if action == "approve":
            selected_resources.update(verification_status="APPROVED", verified_by=request.user, verified_at=timezone.now())
            messages.success(request, "Selected resources approved.")
        elif action == "reject":
            selected_resources.update(verification_status="REJECTED", verified_by=request.user, verified_at=timezone.now())
            messages.success(request, "Selected resources rejected.")
        elif action == "delete":
            for resource in selected_resources:
                resource.file.delete(save=False)
                resource.delete()
            messages.success(request, "Selected resources deleted.")
        return redirect("resources:admin_resource_list")

    paginator = Paginator(resources_qs, 20)
    page_number = request.GET.get("page")
    page_obj = paginator.get_page(page_number)

    context = {
        "page_obj": page_obj,
        "q": q,
        "subject_filter": subject_ids,
        "semester_filter": semesters,
        "status": status,
        "sort": sort,
        "semester_choices": SEMESTER_CHOICES,
        "subjects": Subject.objects.order_by("name"),
        "VERIFICATION_STATUS_CHOICES": [
            ("", "All"),
            ("PENDING", "Pending"),
            ("APPROVED", "Approved"),
            ("REJECTED", "Rejected"),
        ],
    }
    return render(request, "resources/admin_resource_list.html", context)


@staff_member_required
def admin_resource_edit(request, pk):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    resource = get_object_or_404(Resource, pk=pk)

    if request.method == "POST":
        form = AdminResourceForm(request.POST, request.FILES, instance=resource)
        if form.is_valid():
            updated = form.save(commit=False)
            new_subject = form.cleaned_data.get("new_subject_name")
            if new_subject:
                subject_obj, _ = Subject.objects.get_or_create(name=new_subject.strip())
                updated.subject = subject_obj
            updated.save()
            messages.success(request, "Resource updated successfully.")
            return redirect("resources:admin_resource_list")
    else:
        form = AdminResourceForm(instance=resource)

    return render(request, "resources/admin_resource_edit.html", {"form": form, "resource": resource})


@staff_member_required
def admin_report_list(request):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    reports = Report.objects.select_related("resource", "reporter").order_by("-created_at")
    q = request.GET.get("q", "").strip()
    status = request.GET.get("status", "")

    if q:
        reports = reports.filter(
            Q(resource__title__icontains=q)
            | Q(reason__icontains=q)
            | Q(reporter__username__icontains=q)
        )

    if status:
        reports = reports.filter(status=status)

    if request.method == "POST":
        action = request.POST.get("action")
        report_id = request.POST.get("report_id")
        warn_message = request.POST.get("warn_message", "").strip()
        report = get_object_or_404(Report, pk=report_id)

        if action == "approve_report":
            report.status = "RESOLVED"
            report.handled_at = timezone.now()
            report.save()
            messages.success(request, "Report approved.")
        elif action == "reject_report":
            report.status = "REVIEWED"
            report.handled_at = timezone.now()
            report.save()
            messages.success(request, "Report rejected.")
        elif action == "delete_resource":
            resource = report.resource
            resource.file.delete(save=False)
            resource.delete()
            report.status = "RESOLVED"
            report.handled_at = timezone.now()
            report.save()
            messages.success(request, "Resource deleted and report resolved.")
        elif action == "warn_uploader":
            if warn_message:
                Notification.objects.create(
                    user=report.resource.owner,
                    notif_type="REPORT_STATUS",
                    message=f"Warning: {warn_message}",
                    resource=report.resource,
                    report=report,
                )
                report.status = "REVIEWED"
                report.handled_at = timezone.now()
                report.save()
                messages.success(request, "Uploader warned and report updated.")
            else:
                messages.error(request, "Please enter a warning message.")
        return redirect("resources:admin_report_list")

    paginator = Paginator(reports, 20)
    page_number = request.GET.get("page")
    page_obj = paginator.get_page(page_number)

    context = {
        "page_obj": page_obj,
        "q": q,
        "status": status,
        "REPORT_STATUS_CHOICES": [
            ("", "All"),
            ("OPEN", "Open"),
            ("REVIEWED", "Reviewed"),
            ("RESOLVED", "Resolved"),
        ],
    }
    return render(request, "resources/admin_report_list.html", context)


@staff_member_required
def admin_subject_list(request):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    subjects = Subject.objects.order_by("name")
    q = request.GET.get("q", "").strip()
    if q:
        subjects = subjects.filter(
            Q(name__icontains=q)
            | Q(branch__icontains=q)
            | Q(abbreviation__icontains=q)
        )

    context = {
        "subjects": subjects,
        "q": q,
    }
    return render(request, "resources/admin_subject_list.html", context)


@staff_member_required
def admin_subject_add(request):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    if request.method == "POST":
        form = SubjectForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "Subject added successfully.")
            return redirect("resources:admin_subject_list")
    else:
        form = SubjectForm()

    return render(request, "resources/admin_subject_form.html", {"form": form, "title": "Add Subject"})


@staff_member_required
def admin_subject_edit(request, pk):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    subject = get_object_or_404(Subject, pk=pk)
    if request.method == "POST":
        form = SubjectForm(request.POST, instance=subject)
        if form.is_valid():
            form.save()
            messages.success(request, "Subject updated successfully.")
            return redirect("resources:admin_subject_list")
    else:
        form = SubjectForm(instance=subject)

    return render(request, "resources/admin_subject_form.html", {"form": form, "title": "Edit Subject"})


@staff_member_required
def admin_subject_delete(request, pk):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    subject = get_object_or_404(Subject, pk=pk)
    if request.method == "POST":
        subject.delete()
        messages.success(request, "Subject deleted successfully.")
        return redirect("resources:admin_subject_list")

    return render(request, "resources/admin_subject_confirm_delete.html", {"subject": subject})


@staff_member_required
def admin_send_notification(request):
    if not _require_superuser(request):
        return redirect("core:dashboard")

    if request.method == "POST":
        title = request.POST.get("title", "").strip()
        message_text = request.POST.get("message", "").strip()

        if not title or not message_text:
            messages.error(request, "Please enter both a title and a message.")
        else:
            users = User.objects.filter(is_active=True)
            notifications = [
                Notification(
                    user=user,
                    notif_type="REPORT_STATUS",
                    message=f"{title}: {message_text}",
                )
                for user in users
            ]
            Notification.objects.bulk_create(notifications)
            messages.success(request, f"Announcement sent to {users.count()} users.")
            return redirect("resources:admin_send_notification")

    return render(request, "resources/admin_send_notification.html")


# -------------------------------------------------------------------
# Notifications
# -------------------------------------------------------------------
@login_required
def notifications_list(request):
    notifications = request.user.notifications.order_by("-created_at")
    return render(
        request,
        "resources/notifications.html",
        {"notifications": notifications},
    )


@login_required
def notification_mark_read(request, pk):
    notif = get_object_or_404(Notification, pk=pk, user=request.user)
    notif.is_read = True
    notif.save()

    if notif.resource_id:
        return redirect("resources:resource_detail", pk=notif.resource_id)
    return redirect("resources:notifications_list")


# -------------------------------------------------------------------
# Verify resource (staff)
# -------------------------------------------------------------------
@staff_member_required
def verify_resource(request, pk):
    resource = get_object_or_404(Resource, pk=pk)

    if request.method == "POST":
        action = request.POST.get("action")
        note = request.POST.get("note", "").strip()

        if action not in ("APPROVED", "REJECTED"):
            messages.error(request, "Invalid action.")
            return redirect("resources:resource_detail", pk=pk)

        resource.verification_status = action
        resource.verified_by = request.user
        resource.verified_at = timezone.now()
        resource.verification_note = note
        resource.save()

        msg = f"Your resource '{resource.title}' was "
        msg += "approved ✅" if action == "APPROVED" else "rejected ❌"
        if note:
            msg += f" – Note: {note}"

        Notification.objects.create(
            user=resource.owner,
            notif_type="REPORT_STATUS",
            message=msg,
            resource=resource,
        )

        messages.success(request, "Verification status updated.")
        return redirect("resources:resource_detail", pk=pk)

    return redirect("resources:resource_detail", pk=pk)


# -------------------------------------------------------------------
# Leaderboard
# -------------------------------------------------------------------
@login_required
def leaderboard(request):
    """
    Simple leaderboard for top contributors.
    Score = uploads*3 + total_downloads + comments_made
    """
    users = (
        User.objects.select_related("profile").annotate(
            uploads_count=Count("resources", distinct=True),
            total_downloads=Sum("resources__download_count"),
            total_views=Sum("resources__view_count"),
            comments_made=Count("comments", distinct=True),
            ratings_given=Count("ratings", distinct=True),
        )
        .annotate(
            score=(
                F("uploads_count") * 3
                + Coalesce(F("total_downloads"), 0)
                + F("comments_made")
            )
        )
        .order_by("-score")[:20]
    )

    return render(request, "resources/leaderboard.html", {"users": users})


# -------------------------------------------------------------------
# Delete resource (owner or staff only)
# -------------------------------------------------------------------
@login_required
def delete_resource(request, pk):
    resource = get_object_or_404(Resource, pk=pk)

    # Only the owner or staff can delete
    if resource.owner != request.user and not request.user.is_staff:
        messages.error(request, "You do not have permission to delete this resource.")
        return redirect("resources:resource_detail", pk=pk)

    if request.method == "POST":
        title = resource.title
        resource.file.delete(save=False)  # Remove file from disk
        resource.delete()
        messages.success(request, f"Resource '{title}' has been deleted.")
        return redirect("resources:resource_list")

    return render(request, "resources/resource_confirm_delete.html", {"resource": resource})


# -------------------------------------------------------------------
# Mark all notifications as read
# -------------------------------------------------------------------
@login_required
def mark_all_notifications_read(request):
    request.user.notifications.filter(is_read=False).update(is_read=True)
    messages.success(request, "All notifications marked as read.")
    return redirect("resources:notifications_list")
